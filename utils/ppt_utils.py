# -*- coding: utf-8 -*-
"""
python-pptx 기반 프레젠테이션 빌더.

주의(중요):
- Streamlit Community Cloud 런타임에는 LibreOffice가 없는 경우가 대부분이므로,
  PPT 생성은 pptxgenjs가 아니라 python-pptx로 "직접" 슬라이드 XML 요소를 조립한다.
- text_frame.text = "..." 는 서식을 깨뜨리므로 항상 run 단위로 텍스트를 넣는다.
- 색상 테마는 병원 마케팅에 무난한 네이비/골드, 민트, 코랄 3종을 기본 제공한다.
"""

from __future__ import annotations
import io
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


@dataclass
class Theme:
    name: str
    primary: str      # 강조색 (진한 색)
    secondary: str     # 보조 강조색
    bg: str = "FFFFFF"
    text_dark: str = "1A1A1A"
    text_light: str = "6B6B6B"
    card_bg: str = "F4F5F7"


THEMES = {
    "navy_gold": Theme("네이비 & 골드 (신뢰감)", primary="0F2A47", secondary="C9A24B"),
    "mint": Theme("민트 (청결/헬스케어)", primary="0E7C61", secondary="1FB68C"),
    "coral": Theme("코랄 (친근함/여성클리닉)", primary="B5432D", secondary="E8734A"),
}


def _hex(c: str) -> RGBColor:
    return RGBColor.from_string(c)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def _set_fill(shape, color_hex: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color_hex)
    shape.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, size, color_hex,
                  bold=False, align=PP_ALIGN.LEFT, font_name="Malgun Gothic",
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = _hex(color_hex)
    return box


def add_title_slide(prs: Presentation, theme: Theme, title: str, subtitle: str = "",
                     eyebrow: str = ""):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, theme.primary)
    bg.shadow.inherit = False

    # 우측 하단 포인트 원형 (은은한 장식)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(SLIDE_W_IN - 3.2), Inches(SLIDE_H_IN - 3.2),
        Inches(5), Inches(5),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = _hex(theme.secondary)
    circle.fill.transparency = 0
    circle.line.fill.background()
    circle.shadow.inherit = False
    # 투명도 적용 (python-pptx는 transparency 직접 미지원 -> XML 조작)
    _apply_alpha(circle, 18)

    if eyebrow:
        _add_textbox(slide, 0.9, 1.6, 8, 0.5, eyebrow.upper(), 14, theme.secondary, bold=True)
    _add_textbox(slide, 0.9, 2.15, 10.5, 2.2, title, 40, "FFFFFF", bold=True)
    if subtitle:
        _add_textbox(slide, 0.9, 4.1, 9.5, 1.2, subtitle, 18, "E8E8E8")
    return slide


def _apply_alpha(shape, alpha_percent: int):
    """도형 채우기에 투명도를 적용 (alpha_percent: 불투명도 %, 낮을수록 더 투명)"""
    sp = shape.fill.fore_color._xFill
    srgbClr = sp.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_percent * 1000)})
    srgbClr.append(alpha)


def add_section_slide(prs: Presentation, theme: Theme, section_no: str, title: str):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, theme.bg)
    bg.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
    _set_fill(accent, theme.primary)
    accent.shadow.inherit = False

    _add_textbox(slide, 0.9, 2.8, 3, 1, section_no, 20, theme.secondary, bold=True)
    _add_textbox(slide, 0.9, 3.35, 10.5, 1.5, title, 34, theme.text_dark, bold=True)
    return slide


def add_bullet_slide(prs: Presentation, theme: Theme, title: str, bullets: list[str],
                      notes: str | None = None):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, theme.bg)
    bg.shadow.inherit = False

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    _set_fill(top_bar, theme.primary)
    top_bar.shadow.inherit = False

    _add_textbox(slide, 0.7, 0.28, 11.5, 0.7, title, 24, "FFFFFF", bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    card_top = 1.55
    card_h = SLIDE_H_IN - card_top - 0.6
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(card_top),
                                   Inches(SLIDE_W_IN - 1.4), Inches(card_h))
    card.adjustments[0] = 0.03
    _set_fill(card, theme.card_bg)
    card.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(1.1), Inches(card_top + 0.35),
                                    Inches(SLIDE_W_IN - 2.2), Inches(card_h - 0.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(17)
        run.font.name = "Malgun Gothic"
        run.font.color.rgb = _hex(theme.text_dark)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_slide(prs: Presentation, theme: Theme, title: str, image_bytes: bytes,
                     caption: str = "", notes: str | None = None):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, theme.bg)
    bg.shadow.inherit = False

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    _set_fill(top_bar, theme.primary)
    top_bar.shadow.inherit = False
    _add_textbox(slide, 0.7, 0.22, 11.5, 0.6, title, 22, "FFFFFF", bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    img = Image.open(io.BytesIO(image_bytes))
    iw, ih = img.size
    max_w, max_h = SLIDE_W_IN - 2.4, SLIDE_H_IN - 2.4
    ratio = min(max_w / (iw / 96), max_h / (ih / 96))
    disp_w_in = (iw / 96) * ratio
    disp_h_in = (ih / 96) * ratio
    left = (SLIDE_W_IN - disp_w_in) / 2
    top = 1.25 + (max_h - disp_h_in) / 2

    stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(stream, Inches(left), Inches(top), Inches(disp_w_in), Inches(disp_h_in))

    if caption:
        _add_textbox(slide, 0.7, SLIDE_H_IN - 0.55, SLIDE_W_IN - 1.4, 0.4, caption, 12,
                     theme.text_light, align=PP_ALIGN.CENTER)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_before_after_slide(prs: Presentation, theme: Theme, title: str,
                            before_bytes: bytes, after_bytes: bytes,
                            before_label="Before", after_label="After"):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, theme.bg)
    bg.shadow.inherit = False

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    _set_fill(top_bar, theme.primary)
    top_bar.shadow.inherit = False
    _add_textbox(slide, 0.7, 0.22, 11.5, 0.6, title, 22, "FFFFFF", bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    col_w = (SLIDE_W_IN - 1.4 - 0.4) / 2
    img_h = SLIDE_H_IN - 2.3

    for idx, (label, img_bytes) in enumerate([(before_label, before_bytes), (after_label, after_bytes)]):
        left = 0.7 + idx * (col_w + 0.4)
        stream = io.BytesIO(img_bytes)
        pic = slide.shapes.add_picture(stream, Inches(left), Inches(1.35), height=Inches(img_h))
        # 폭이 컬럼을 넘으면 폭 기준으로 재조정
        if pic.width > Inches(col_w):
            ratio = Inches(col_w) / pic.width
            pic.width = Inches(col_w)
            pic.height = int(pic.height * ratio)
        _add_textbox(slide, left, 1.35 + img_h + 0.1, col_w, 0.5, label, 16, theme.primary,
                     bold=True, align=PP_ALIGN.CENTER)
    return slide


def presentation_to_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
